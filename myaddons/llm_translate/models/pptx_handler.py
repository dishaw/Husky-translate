"""PowerPoint (.ppt/.pptx) parsing and reconstruction utilities.

Uses python-pptx for .pptx files.
For legacy .ppt files, uses LibreOffice headless to convert to .pptx first.
"""

import base64
import io
import json
import logging
import os
import re
import subprocess
import tempfile

_logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
except ImportError:
    Presentation = None
    qn = None
    _logger.warning("python-pptx not installed. PPTX support unavailable.")


def estimate_tokens(text):
    """Estimate token count for a text string."""
    if not text:
        return 0
    cjk_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff'
                    or '\u3400' <= c <= '\u4dbf'
                    or '\uf900' <= c <= '\ufaff'
                    or '\u3000' <= c <= '\u303f'
                    or '\u3040' <= c <= '\u309f'
                    or '\u30a0' <= c <= '\u30ff'
                    or '\uac00' <= c <= '\ud7af')
    ascii_count = len(text) - cjk_count
    return int(cjk_count / 1.5 + ascii_count / 4)


def split_long_paragraph(text, max_tokens=2000):
    """Split a paragraph that exceeds max_tokens into sentence-level chunks."""
    if estimate_tokens(text) <= max_tokens:
        return [text]

    sentence_pattern = r'(?<=[。.！!？?；;\n])'
    sentences = re.split(sentence_pattern, text)
    sentences = [s for s in sentences if s.strip()]

    if not sentences:
        return [text]

    chunks = []
    current_chunk = ""

    for sentence in sentences:
        test_chunk = current_chunk + sentence
        if estimate_tokens(test_chunk) > max_tokens and current_chunk:
            chunks.append(current_chunk)
            current_chunk = sentence
        else:
            current_chunk = test_chunk

    if current_chunk:
        chunks.append(current_chunk)

    return chunks if chunks else [text]


def _convert_ppt_to_pptx_via_libreoffice(file_content):
    """Convert legacy .ppt to .pptx using LibreOffice headless.

    Args:
        file_content: Binary content of the .ppt file.

    Returns:
        bytes: Binary content of the converted .pptx file.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, "input.ppt")
        with open(input_path, "wb") as f:
            f.write(file_content)

        cmd = [
            "libreoffice", "--headless", "--convert-to", "pptx",
            "--outdir", tmpdir, input_path
        ]
        try:
            subprocess.run(cmd, check=True, timeout=120,
                           capture_output=True, text=True)
        except FileNotFoundError:
            raise RuntimeError(
                "LibreOffice is not installed or not in PATH. "
                "It is required to convert .ppt files to .pptx."
            )
        except subprocess.TimeoutExpired:
            raise RuntimeError("LibreOffice conversion timed out after 120 seconds.")
        except subprocess.CalledProcessError as e:
            raise RuntimeError(f"LibreOffice conversion failed: {e.stderr}")

        output_path = os.path.join(tmpdir, "input.pptx")
        if not os.path.exists(output_path):
            raise RuntimeError("LibreOffice conversion produced no output file.")

        with open(output_path, "rb") as f:
            return f.read()


def _extract_run_info(run):
    """Extract formatting info from a pptx run."""
    info = {
        "text": run.text,
        "bold": False,
        "italic": False,
        "font_size": None,
        "font_name": None,
        "color": None,
    }
    font = run.font
    if font.bold:
        info["bold"] = True
    if font.italic:
        info["italic"] = True
    if font.size:
        info["font_size"] = font.size.pt
    if font.name:
        info["font_name"] = font.name
    if font.color:
        try:
            if font.color.rgb:
                info["color"] = str(font.color.rgb)
        except (AttributeError, TypeError):
            pass
    return info


def _get_alignment_str(alignment):
    """Convert PP_ALIGN enum to string."""
    if alignment is None:
        return None
    align_map = {
        PP_ALIGN.LEFT: "LEFT",
        PP_ALIGN.CENTER: "CENTER",
        PP_ALIGN.RIGHT: "RIGHT",
        PP_ALIGN.JUSTIFY: "JUSTIFY",
    }
    return align_map.get(alignment)


def _extract_shape_text(shape):
    """Extract text paragraphs from a shape's text_frame.

    Returns:
        list[dict]: paragraph dicts with text, runs, alignment, etc.
    """
    if not shape.has_text_frame:
        return []

    paragraphs = []
    tf = shape.text_frame
    for para in tf.paragraphs:
        text = para.text
        runs = []
        for run in para.runs:
            runs.append(_extract_run_info(run))

        alignment = _get_alignment_str(para.alignment)
        # Determine bold from first run
        bold = runs[0]["bold"] if runs else False
        font_size = runs[0].get("font_size") if runs else None

        paragraphs.append({
            "text": text,
            "runs": runs,
            "alignment": alignment,
            "bold": bold,
            "font_size": font_size,
            "is_empty": not text.strip(),
        })

    return paragraphs


def _extract_table_text(table):
    """Extract text from a table shape.

    Returns:
        list[dict]: One dict per row with cells joined by [CELL] separator.
    """
    rows_data = []
    for row_idx, row in enumerate(table.rows):
        cell_texts = []
        for cell in row.cells:
            cell_texts.append(cell.text.strip())
        combined = " [CELL] ".join(cell_texts)
        rows_data.append({
            "text": combined,
            "is_table_row": True,
            "is_empty": not combined.strip(),
            "row_index": row_idx,
            "num_cells": len(row.cells),
        })
    return rows_data


def extract_paragraphs_from_pptx(file_content):
    """Extract text from a .pptx file organized by slides and shapes.

    Args:
        file_content: Binary content of the .pptx file.

    Returns:
        dict: {
            "paragraphs": list[dict],  # flat list of extractable text segments
        }

    Each paragraph dict contains:
        - text (str)
        - slide_index (int): 0-based slide number
        - shape_index (int): shape index within slide
        - para_index (int): global paragraph index
        - shape_name (str): shape name
        - style (str): 'Title', 'Subtitle', 'Body', 'Table', etc.
        - alignment, bold, font_size, runs
        - is_empty (bool)
        - is_table_row (bool)
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for .pptx support. "
            "Install with: pip install python-pptx"
        )

    prs = Presentation(io.BytesIO(file_content))
    all_paragraphs = []
    para_index = 0

    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = shape.name or f"Shape {shape_idx}"

            # Determine style hint from shape type/placeholder
            style = "Body"
            if shape.has_text_frame:
                try:
                    ph_fmt = shape.placeholder_format
                    if ph_fmt is not None:
                        ph_idx = ph_fmt.idx
                        if ph_idx == 0:
                            style = "Title"
                        elif ph_idx == 1:
                            style = "Subtitle"
                except ValueError:
                    # shape.placeholder_format raises ValueError for non-placeholder shapes
                    pass

            # Handle tables
            if shape.has_table:
                table_rows = _extract_table_text(shape.table)
                for row_data in table_rows:
                    row_data.update({
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_name": shape_name,
                        "para_index": para_index,
                        "style": "Table",
                    })
                    all_paragraphs.append(row_data)
                    para_index += 1
                continue

            # Handle text shapes
            if shape.has_text_frame:
                shape_paras = _extract_shape_text(shape)
                for sp in shape_paras:
                    sp.update({
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_name": shape_name,
                        "para_index": para_index,
                        "style": style,
                        "is_table_row": False,
                    })
                    all_paragraphs.append(sp)
                    para_index += 1

        # Also extract from notes slide if present
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                all_paragraphs.append({
                    "text": notes_text,
                    "slide_index": slide_idx,
                    "shape_index": -1,
                    "shape_name": "Notes",
                    "para_index": para_index,
                    "style": "Notes",
                    "alignment": None,
                    "bold": False,
                    "font_size": None,
                    "runs": [],
                    "is_empty": False,
                    "is_table_row": False,
                })
                para_index += 1

    return {"paragraphs": all_paragraphs}


def extract_paragraphs_from_ppt(file_content):
    """Extract paragraphs from a legacy .ppt file.

    Converts to .pptx via LibreOffice, then processes as .pptx.

    Args:
        file_content: Binary content of the .ppt file.

    Returns:
        dict: Same format as extract_paragraphs_from_pptx.
    """
    pptx_content = _convert_ppt_to_pptx_via_libreoffice(file_content)
    return extract_paragraphs_from_pptx(pptx_content)


def rebuild_pptx_from_original(original_content, paragraphs_data):
    """Rebuild a .pptx file by replacing text in-place, preserving formatting.

    Args:
        original_content: Binary content of the original .pptx.
        paragraphs_data: dict with:
            - paragraphs: list[dict] each with 'translated_text' and 'style_metadata'

    Returns:
        bytes: Binary content of the translated .pptx.
    """
    if Presentation is None:
        raise ImportError(
            "python-pptx is required for .pptx support. "
            "Install with: pip install python-pptx"
        )

    prs = Presentation(io.BytesIO(original_content))

    # Build a lookup: (slide_index, shape_index, local_para_index) -> translated_text
    # For tables: (slide_index, shape_index, row_index) -> translated_text
    text_lookup = {}
    table_lookup = {}

    paras = paragraphs_data.get("paragraphs", [])
    for pdata in paras:
        meta = pdata.get("style_metadata", {})
        translated = pdata.get("translated_text", "")
        slide_idx = meta.get("slide_index")
        shape_idx = meta.get("shape_index")

        if slide_idx is None or shape_idx is None:
            continue

        if meta.get("is_table_row"):
            row_idx = meta.get("row_index", 0)
            table_lookup[(slide_idx, shape_idx, row_idx)] = translated
        else:
            para_idx = meta.get("para_index")
            if para_idx is not None:
                text_lookup[(slide_idx, shape_idx, para_idx)] = translated

    # We also need to track per-shape local paragraph indices
    # Rebuild local para index by iterating the same way we extracted
    local_para_counter = {}  # (slide_idx, shape_idx) -> next local index

    para_index = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # Handle tables
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    key = (slide_idx, shape_idx, row_idx)
                    translated = table_lookup.get(key)
                    if translated:
                        cells_text = translated.split(" [CELL] ")
                        for cell_idx, cell in enumerate(row.cells):
                            if cell_idx < len(cells_text):
                                _replace_text_frame_content(
                                    cell.text_frame, cells_text[cell_idx]
                                )
                    para_index += 1
                continue

            # Handle text shapes
            if shape.has_text_frame:
                tf = shape.text_frame
                has_translation = False
                for local_para_idx, para in enumerate(tf.paragraphs):
                    key = (slide_idx, shape_idx, para_index)
                    translated = text_lookup.get(key)
                    if translated and translated.strip():
                        _replace_paragraph_text(para, translated)
                        has_translation = True
                    para_index += 1
                if has_translation:
                    _auto_fit_text_frame(shape)
                continue

            para_index += 0  # non-text shapes don't increment

        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                key = (slide_idx, -1, para_index)
                translated = text_lookup.get(key)
                if translated and translated.strip():
                    _replace_text_frame_content(
                        slide.notes_slide.notes_text_frame, translated
                    )
                para_index += 1

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _normalize_line_spacing(paragraph):
    """Normalize paragraph line spacing to single (100%) to prevent
    line height being much larger than font size.

    Removes any fixed line spacing (spcPts) and sets proportional
    spacing to 100% (single line).
    """
    pPr = paragraph._p.find(qn('a:pPr'))
    if pPr is None:
        return
    lnSpc = pPr.find(qn('a:lnSpc'))
    if lnSpc is not None:
        pPr.remove(lnSpc)
    # Also remove excessive space before/after
    for tag in ('a:spcBef', 'a:spcAft'):
        elem = pPr.find(qn(tag))
        if elem is not None:
            pPr.remove(elem)


def _auto_fit_text_frame(shape):
    """Configure text frame: keep original box size/position, enable word wrap.

    - Uses noAutofit so the box stays exactly the same size as original.
    - Enables word wrap so long translated text wraps instead of overflowing.
    - Keeps original font size (no shrinking).
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Enable word-wrap so long text wraps within the fixed-size box
    tf.word_wrap = True
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        # Remove any existing autofit setting
        for child_tag in ('a:noAutofit', 'a:spAutoFit', 'a:normAutofit'):
            child = bodyPr.find(qn(child_tag))
            if child is not None:
                bodyPr.remove(child)
        # Use noAutofit: box size stays fixed, text wraps within it
        from lxml import etree
        etree.SubElement(bodyPr, qn('a:noAutofit'))


def _replace_paragraph_text(paragraph, new_text):
    """Replace the text of a paragraph while preserving first-run formatting.

    If the paragraph has runs, keeps the first run's formatting and sets its text
    to new_text, removing subsequent runs. If no runs, sets the paragraph text directly.
    Also normalizes line spacing to prevent height mismatch.
    """
    if not paragraph.runs:
        # No runs - try to set text via a new run
        paragraph.text = new_text
        _normalize_line_spacing(paragraph)
        return

    # Preserve first run formatting, set its text
    first_run = paragraph.runs[0]
    first_run.text = new_text

    # Remove subsequent runs by clearing their XML elements
    for run in paragraph.runs[1:]:
        run._r.getparent().remove(run._r)

    # Normalize line spacing to prevent 3x height issue
    _normalize_line_spacing(paragraph)


def _replace_text_frame_content(text_frame, new_text):
    """Replace all text in a text_frame with new_text, preserving first paragraph formatting."""
    if not text_frame.paragraphs:
        return

    # Set first paragraph
    first_para = text_frame.paragraphs[0]
    _replace_paragraph_text(first_para, new_text)

    # Remove extra paragraphs
    from lxml import etree
    body = text_frame._txBody
    p_elements = body.findall(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}p'
    )
    for p_elem in p_elements[1:]:
        body.remove(p_elem)
